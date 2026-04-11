"""
DevBot — Custom AI Assistant
Gemini-powered chatbot with multi-session management, knowledge base injection,
dynamic role assignment, and a polished Streamlit UI.

Usage:
    1. Set GEMINI_API_KEY in your environment (or create config.py with GEMINI_API_KEY = "...")
    2. Run:  streamlit run devbot.py
"""

from __future__ import annotations

import json
import logging
import os
import time
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Generator

import streamlit as st
from google import genai
from google.genai import types
from google.genai.errors import ClientError

# --- NEW IMPORTS FOR FILE SUPPORT ---
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation

# ─────────────────────────────────────────────────────────────────────────────
# Configuration constants
# ─────────────────────────────────────────────────────────────────────────────

KB_DIR            = Path("knowledge")   # Directory for knowledge-base files
LOG_FILE          = Path("devbot.log")
SESSIONS_DIR      = Path("sessions")

GEMINI_MODEL      = "gemini-2.5-flash"
MAX_HISTORY       = 20          
MAX_OUTPUT_TOKENS = 4096        
TEMPERATURE       = 0.65        
MAX_RETRIES       = 3
RETRY_BASE_DELAY  = 2           

SYSTEM_PROMPT_TEMPLATE = """\
You are DevBot, an expert {role}.
You provide high-quality, accurate, and professional responses.
Always prefer modern best practices, flag potential issues proactively,
and suggest improvements when relevant.

When answering:
- Use code blocks or structured formatting where appropriate.
- Prefer concise explanations backed by working examples.
- Mention trade-offs when multiple valid approaches exist.
- If you are unsure about something, say so rather than guessing.

{knowledge_base}
"""

# ─────────────────────────────────────────────────────────────────────────────
# Logging
# ─────────────────────────────────────────────────────────────────────────────

def _setup_logger() -> logging.Logger:
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(name)s — %(message)s",
        handlers=[
            logging.FileHandler(LOG_FILE, encoding="utf-8"),
            logging.StreamHandler(),
        ],
    )
    return logging.getLogger("devbot")

logger = _setup_logger()

# ─────────────────────────────────────────────────────────────────────────────
# Session persistence
# ─────────────────────────────────────────────────────────────────────────────

def _session_path(session_id: str) -> Path:
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    return SESSIONS_DIR / f"{session_id}.json"

def save_session(session_id: str, title: str, messages: list[dict]) -> None:
    data = {
        "id":         session_id,
        "title":      title,
        "updated_at": datetime.now(timezone.utc).isoformat(),
        "messages":   messages,
    }
    try:
        _session_path(session_id).write_text(
            json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
    except OSError as exc:
        logger.error("Failed to save session %s: %s", session_id, exc)

def load_session(session_id: str) -> dict | None:
    path = _session_path(session_id)
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        logger.error("Failed to load session %s: %s", session_id, exc)
        return None

def list_sessions() -> list[dict]:
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    sessions: list[dict] = []
    for path in SESSIONS_DIR.glob("*.json"):
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
            sessions.append({
                "id":         data.get("id", path.stem),
                "title":      data.get("title", "Untitled"),
                "updated_at": data.get("updated_at", ""),
                "msg_count":  len(data.get("messages", [])),
            })
        except (OSError, json.JSONDecodeError):
            continue
    return sorted(sessions, key=lambda s: s["updated_at"], reverse=True)

def delete_session(session_id: str) -> None:
    path = _session_path(session_id)
    if path.exists():
        path.unlink()
        logger.info("Deleted session %s", session_id)

def derive_title(messages: list[dict]) -> str:
    for m in messages:
        if m["role"] == "user":
            text = m["content"].strip().replace("\n", " ")
            return text[:60] + ("..." if len(text) > 60 else "")
    return "New Chat"

# ─────────────────────────────────────────────────────────────────────────────
# Knowledge base (Enhanced for PDF, Docx, Xlsx, Pptx)
# ─────────────────────────────────────────────────────────────────────────────

def _extract_text(path: Path) -> str:
    """Helper to extract text from various file formats."""
    ext = path.suffix.lower()
    
    try:
        if ext == ".txt":
            return path.read_text(encoding="utf-8").strip()
        
        elif ext == ".json":
            raw = path.read_text(encoding="utf-8").strip()
            try:
                return json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
            except json.JSONDecodeError:
                return raw

        elif ext == ".pdf":
            text = ""
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()

        elif ext == ".docx":
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs]).strip()

        elif ext == ".xlsx":
            # Convert all sheets to a string representation
            df_dict = pd.read_excel(path, sheet_name=None)
            sheets_text = []
            for sheet_name, df in df_dict.items():
                sheets_text.append(f"Sheet: {sheet_name}\n{df.to_string()}")
            return "\n\n".join(sheets_text).strip()

        elif ext == ".pptx":
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text).strip()

    except Exception as e:
        logger.error("Error extracting text from %s: %s", path, e)
        return f"[Error reading file {path.name}: {e}]"

    return ""

@st.cache_data(show_spinner=False)
def load_knowledge_base(kb_dir: str = str(KB_DIR)) -> tuple[str, list[str]]:
    dir_path = Path(kb_dir)
    if not dir_path.is_dir():
        return "", []

    sections: list[str] = []
    names:    list[str] = []
    
    # Supported extensions
    valid_extensions = {".txt", ".json", ".pdf", ".docx", ".xlsx", ".pptx"}

    for path in sorted(dir_path.glob("*")):
        if path.suffix.lower() not in valid_extensions:
            continue
        
        content = _extract_text(path)
        if content:
            sections.append(f"### [{path.name}]\n{content}")
            names.append(path.name)

    logger.info("Knowledge base: %d file(s) loaded -> %s", len(names), names)
    return "\n\n---\n\n".join(sections), names

# ─────────────────────────────────────────────────────────────────────────────
# Gemini API helpers
# ─────────────────────────────────────────────────────────────────────────────

@st.cache_resource(show_spinner=False)
def get_client(api_key: str) -> genai.Client:
    return genai.Client(api_key=api_key)

def build_system_prompt(knowledge_base: str, role: str) -> str:
    kb_section = (
        f"## Knowledge Base\n\n{knowledge_base}"
        if knowledge_base
        else "No additional knowledge base loaded."
    )
    return SYSTEM_PROMPT_TEMPLATE.format(role=role, knowledge_base=kb_section)

def _to_contents(messages: list[dict]) -> list[types.Content]:
    role_map = {"user": "user", "assistant": "model"}
    windowed = messages[-MAX_HISTORY:]
    return [
        types.Content(
            role=role_map.get(m["role"], "user"),
            parts=[types.Part(text=m["content"])],
        )
        for m in windowed
    ]

def stream_response(
    client: genai.Client,
    system_prompt: str,
    messages: list[dict],
) -> Generator[str, None, None]:
    contents = _to_contents(messages)
    cfg = types.GenerateContentConfig(
        system_instruction=system_prompt,
        temperature=TEMPERATURE,
        max_output_tokens=MAX_OUTPUT_TOKENS,
    )

    for attempt in range(MAX_RETRIES + 1):
        try:
            for chunk in client.models.generate_content_stream(
                model=GEMINI_MODEL, contents=contents, config=cfg
            ):
                if chunk.text:
                    yield chunk.text
            return 

        except ClientError as exc:
            logger.warning("ClientError (attempt %d/%d): %s", attempt + 1, MAX_RETRIES + 1, exc)
            if attempt < MAX_RETRIES:
                time.sleep(RETRY_BASE_DELAY ** attempt)
            else:
                yield "API quota exceeded or request rejected. Please wait a moment."
                return
        except Exception as exc:
            logger.error("Unexpected error (attempt %d/%d): %s", attempt + 1, MAX_RETRIES + 1, exc)
            if attempt < MAX_RETRIES:
                time.sleep(RETRY_BASE_DELAY ** attempt)
            else:
                yield f"Unexpected error: {exc}"
                return

# ─────────────────────────────────────────────────────────────────────────────
# Session state helpers
# ─────────────────────────────────────────────────────────────────────────────

def _new_session_id() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S_") + uuid.uuid4().hex[:6]

def _init_state() -> None:
    defaults: dict = {
        "session_id":      _new_session_id(),
        "messages":        [],
        "viewing_session": None,
        "response_cache":  {},
        "bot_role":        "Full-Stack Developer and Software Architect",
    }
    for key, value in defaults.items():
        st.session_state.setdefault(key, value)

def _flush_current_session() -> None:
    sid   = st.session_state.session_id
    msgs  = st.session_state.messages
    title = derive_title(msgs)
    save_session(sid, title, msgs)

def _start_new_chat() -> None:
    st.session_state.session_id      = _new_session_id()
    st.session_state.messages        = []
    st.session_state.viewing_session = None

def _switch_to_session(session_id: str) -> None:
    data = load_session(session_id)
    if data:
        st.session_state.messages        = data.get("messages", [])
        st.session_state.viewing_session = session_id
    else:
        st.warning("Session not found or could not be loaded.")

def _cache_get(prompt: str) -> str | None:
    return st.session_state.response_cache.get(prompt)

def _cache_set(prompt: str, response: str) -> None:
    cache = st.session_state.response_cache
    if len(cache) >= 50:
        oldest = next(iter(cache))
        del cache[oldest]
    cache[prompt] = response

def handle_user_message(prompt: str, client: genai.Client, system_prompt: str) -> None:
    cached = _cache_get(prompt)
    if cached:
        with st.chat_message("assistant"):
            st.markdown(cached)
        st.session_state.messages.append({"role": "assistant", "content": cached})
        return

    full_response = ""
    with st.chat_message("assistant"):
        placeholder = st.empty()
        for chunk in stream_response(client, system_prompt, st.session_state.messages):
            full_response += chunk
            placeholder.markdown(full_response + " |")
        placeholder.markdown(full_response)

    if full_response:
        st.session_state.messages.append({"role": "assistant", "content": full_response})
        _cache_set(prompt, full_response)
        _flush_current_session()

def resolve_api_key() -> str:
    try:
        import config
        key = getattr(config, "GEMINI_API_KEY", None)
        if key: return key
    except ImportError: pass
    return os.getenv("GEMINI_API_KEY", "")

def _fmt_date(iso: str) -> str:
    try:
        dt   = datetime.fromisoformat(iso).replace(tzinfo=timezone.utc)
        diff = (datetime.now(timezone.utc) - dt).total_seconds()
        if diff < 60:    return "just now"
        if diff < 3600:  return f"{int(diff // 60)}m ago"
        if diff < 86400: return f"{int(diff // 3600)}h ago"
        return dt.strftime("%b %d")
    except Exception: return ""

# ─────────────────────────────────────────────────────────────────────────────
# UI — Sidebar
# ─────────────────────────────────────────────────────────────────────────────

def render_sidebar() -> None:
    with st.sidebar:
        st.markdown(
            """
            <div style="display:flex;align-items:center;gap:10px;padding:4px 0 12px">
                <span style="font-size:1.9rem">&#x1F916;</span>
                <div>
                    <div style="font-size:1.25rem;font-weight:700;letter-spacing:-0.5px;line-height:1.1">
                        DevBot
                    </div>
                    <div style="font-size:0.72rem;opacity:0.55">Custom AI Assistant</div>
                </div>
            </div>
            """, unsafe_allow_html=True,
        )
        st.caption("Powered by Gemini 2.5 Flash")
        st.divider()

        st.markdown("#### Bot Persona")
        st.session_state.bot_role = st.text_input(
            "Define the AI's role:", 
            value=st.session_state.bot_role
        )
        st.divider()

        if st.button("+ New Chat", use_container_width=True, type="primary"):
            _start_new_chat()
            st.rerun()

        st.markdown("#### Recent Chats")
        sessions   = list_sessions()
        active_id  = st.session_state.session_id
        viewing_id = st.session_state.get("viewing_session")

        if not sessions:
            st.caption("No saved chats yet.")
        else:
            for s in sessions:
                sid = s["id"]
                is_live = sid == active_id and not viewing_id
                is_view = sid == viewing_id
                icon  = "> " if is_live else ("^ " if is_view else "")
                label = f"{icon}{s['title']}"
                
                col_btn, col_del = st.columns([5, 1])
                with col_btn:
                    if st.button(label, key=f"sess_{sid}", use_container_width=True, 
                                 type="primary" if (is_live or is_view) else "secondary"):
                        if not (is_live):
                            _switch_to_session(sid)
                        st.rerun()
                with col_del:
                    if st.button("X", key=f"del_{sid}"):
                        delete_session(sid)
                        if sid in (active_id, viewing_id): _start_new_chat()
                        st.rerun()

        st.divider()

        with st.expander("Knowledge Base", expanded=False):
            st.caption("Upload files to inject into the system prompt.")
            # UPDATED: Added more extensions to the uploader
            uploaded_files = st.file_uploader(
                "Upload files",
                type=["txt", "json", "pdf", "docx", "xlsx", "pptx"],
                accept_multiple_files=True,
                label_visibility="collapsed",
            )
            if uploaded_files:
                KB_DIR.mkdir(parents=True, exist_ok=True)
                saved = []
                for uf in uploaded_files:
                    dest = KB_DIR / uf.name
                    dest.write_bytes(uf.read())
                    saved.append(uf.name)
                load_knowledge_base.clear() 
                st.success(f"Saved {len(saved)} file(s)")

            _, kb_files = load_knowledge_base()
            if kb_files:
                for name in kb_files: st.caption(f"* {name}")
            else:
                st.caption("No knowledge files loaded.")

        with st.expander("Settings", expanded=False):
            st.caption(f"Model: {GEMINI_MODEL}")
            st.caption(f"Max output tokens: {MAX_OUTPUT_TOKENS}")
            st.caption(f"Temp: {TEMPERATURE}")

# ─────────────────────────────────────────────────────────────────────────────
# UI — Main chat area
# ─────────────────────────────────────────────────────────────────────────────

def render_chat_area(client: genai.Client, system_prompt: str) -> None:
    viewing      = st.session_state.get("viewing_session")
    is_readonly  = viewing is not None
    messages     = st.session_state.messages

    title = derive_title(messages) if messages else "New Chat"
    hcol, bcol = st.columns([5, 1])
    with hcol: st.markdown(f"### {title}")
    with bcol:
        if is_readonly:
            st.markdown('<span style="background:#f0ad4e;color:#000;padding:3px 12px;border-radius:20px;font-size:0.75rem;font-weight:600">Read-only</span>', unsafe_allow_html=True)
            if st.button("Resume", use_container_width=True):
                data = load_session(st.session_state.session_id)
                st.session_state.messages = data["messages"] if data else []
                st.session_state.viewing_session = None
                st.rerun()

    st.divider()

    if not messages:
        st.markdown(f"""<div style="text-align:center;padding:5rem 1rem;opacity:0.4">
                <div style="font-size:4rem">&#x1F916;</div>
                <p style="font-size:1.2rem;margin-top:0.6rem;font-weight:500">I am configured as a <b>{st.session_state.bot_role}</b>.</p>
            </div>""", unsafe_allow_html=True)
    else:
        for msg in messages:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

    if is_readonly:
        st.info("Viewing a past session. Click **Resume** to continue.", icon="ℹ️")
        return

    if prompt := st.chat_input("Ask anything..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        handle_user_message(prompt, client, system_prompt)

def main() -> None:
    st.set_page_config(page_title="DevBot - AI Assistant", page_icon="robot", layout="wide")
    _init_state()

    api_key = resolve_api_key()
    if not api_key:
        st.error("Gemini API key not found.")
        st.stop()

    client            = get_client(api_key)
    knowledge_base, _ = load_knowledge_base()
    system_prompt     = build_system_prompt(knowledge_base, st.session_state.bot_role)

    render_sidebar()
    render_chat_area(client, system_prompt)

if __name__ == "__main__":
    main()
